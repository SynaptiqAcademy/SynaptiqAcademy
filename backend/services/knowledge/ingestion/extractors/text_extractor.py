"""Extractors for plain text, Markdown, HTML, CSV, and PowerPoint."""
from __future__ import annotations

import csv
import io
import re
import zipfile

from services.knowledge.ingestion.extractors.base import DocumentExtractor, ExtractedDocument
from services.knowledge.models import DocumentMetadata

_MD_HEADING = re.compile(r"^(#{1,6})\s+(.+)$", re.MULTILINE)
_HTML_TAG = re.compile(r"<[^>]+>")
_HTML_ENTITY = re.compile(r"&[a-zA-Z]+;|&#\d+;")


def _simple_sections(text: str, heading_re=None) -> list[dict]:
    """Generic section splitter from detected headings."""
    if heading_re is None:
        return [{"heading": "", "text": text.strip(), "page": None}]
    sections: list[dict] = []
    last_end = 0
    last_heading = ""
    for m in heading_re.finditer(text):
        chunk = text[last_end:m.start()].strip()
        if chunk:
            sections.append({"heading": last_heading, "text": chunk, "page": None})
        last_heading = m.group(0).lstrip("#").strip() if m.lastindex and m.lastindex >= 1 else m.group(0)
        last_end = m.end()
    tail = text[last_end:].strip()
    if tail:
        sections.append({"heading": last_heading, "text": tail, "page": None})
    return sections or [{"heading": "", "text": text, "page": None}]


class TextExtractor(DocumentExtractor):
    """Plain text (.txt) extractor."""

    @property
    def supported_types(self) -> list[str]:
        return ["txt"]

    def extract(self, content: bytes, filename: str = "") -> ExtractedDocument:
        text = content.decode("utf-8", errors="replace")
        text = re.sub(r"\r\n?", "\n", text)
        return ExtractedDocument(
            text=text,
            sections=[{"heading": "", "text": text, "page": None}],
            metadata=DocumentMetadata(title=filename),
            extraction_method="plain",
        )


class MarkdownExtractor(DocumentExtractor):
    """Markdown (.md) extractor — preserves heading structure."""

    @property
    def supported_types(self) -> list[str]:
        return ["md", "markdown"]

    def extract(self, content: bytes, filename: str = "") -> ExtractedDocument:
        text = content.decode("utf-8", errors="replace")
        # Extract title from first H1
        title_m = re.search(r"^#\s+(.+)$", text, re.MULTILINE)
        title = title_m.group(1).strip() if title_m else filename
        sections = _simple_sections(text, _MD_HEADING)
        # Strip markdown formatting for clean text
        clean = re.sub(r"[*_`~]", "", text)
        clean = re.sub(r"!\[.*?\]\(.*?\)", "", clean)
        clean = re.sub(r"\[([^\]]+)\]\([^\)]+\)", r"\1", clean)
        return ExtractedDocument(
            text=clean,
            sections=sections,
            metadata=DocumentMetadata(title=title),
            extraction_method="markdown",
        )


class HTMLExtractor(DocumentExtractor):
    """HTML extractor — strips tags."""

    @property
    def supported_types(self) -> list[str]:
        return ["html", "htm"]

    def extract(self, content: bytes, filename: str = "") -> ExtractedDocument:
        raw = content.decode("utf-8", errors="replace")
        # Extract title
        title_m = re.search(r"<title[^>]*>([^<]+)</title>", raw, re.IGNORECASE)
        title = title_m.group(1).strip() if title_m else filename
        # Strip scripts and styles
        raw = re.sub(r"<(script|style)[^>]*>.*?</(script|style)>", "", raw,
                     flags=re.DOTALL | re.IGNORECASE)
        text = _HTML_TAG.sub(" ", raw)
        text = _HTML_ENTITY.sub(" ", text)
        text = re.sub(r"\s+", " ", text).strip()
        return ExtractedDocument(
            text=text,
            sections=[{"heading": "", "text": text, "page": None}],
            metadata=DocumentMetadata(title=title),
            extraction_method="html",
        )


class CSVExtractor(DocumentExtractor):
    """CSV extractor — converts to structured text."""

    @property
    def supported_types(self) -> list[str]:
        return ["csv"]

    def extract(self, content: bytes, filename: str = "") -> ExtractedDocument:
        try:
            text_content = content.decode("utf-8", errors="replace")
            reader = csv.reader(io.StringIO(text_content))
            rows = list(reader)
        except Exception:
            rows = []

        if not rows:
            return ExtractedDocument(
                text="", sections=[], metadata=DocumentMetadata(title=filename),
                extraction_method="csv",
            )

        header = rows[0] if rows else []
        lines: list[str] = []
        if header:
            lines.append(" | ".join(header))
            lines.append("-" * 40)
        for row in rows[1:101]:  # limit to 100 data rows for indexing
            lines.append(" | ".join(str(v) for v in row))

        text = "\n".join(lines)
        return ExtractedDocument(
            text=text,
            sections=[{"heading": "Table Data", "text": text, "page": None}],
            metadata=DocumentMetadata(title=filename, keywords=header[:10]),
            extraction_method="csv",
        )


class PowerPointExtractor(DocumentExtractor):
    """PPTX extractor — extracts slide text using python-pptx or XML fallback."""

    @property
    def supported_types(self) -> list[str]:
        return ["pptx", "ppt"]

    def extract(self, content: bytes, filename: str = "") -> ExtractedDocument:
        text, sections = self._try_pptx(content)
        if not text.strip():
            text, sections = self._xml_fallback(content)
        return ExtractedDocument(
            text=text,
            sections=sections or [{"heading": "", "text": text, "page": None}],
            metadata=DocumentMetadata(title=filename),
            extraction_method="pptx",
        )

    def _try_pptx(self, content: bytes) -> tuple[str, list[dict]]:
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            sections: list[dict] = []
            for i, slide in enumerate(prs.slides, 1):
                slide_texts: list[str] = []
                title = ""
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    t = shape.text_frame.text.strip()
                    if not t:
                        continue
                    if shape.name.lower().startswith("title"):
                        title = t
                    else:
                        slide_texts.append(t)
                if slide_texts or title:
                    sections.append({
                        "heading": title or f"Slide {i}",
                        "text": "\n".join(slide_texts),
                        "page": i,
                    })
            full = "\n\n".join(
                f"{s['heading']}\n{s['text']}" if s["heading"] else s["text"]
                for s in sections
            )
            return full, sections
        except ImportError:
            return "", []
        except Exception:
            return "", []

    def _xml_fallback(self, content: bytes) -> tuple[str, list[dict]]:
        try:
            with zipfile.ZipFile(io.BytesIO(content)) as z:
                slides = sorted(
                    [n for n in z.namelist() if re.match(r"ppt/slides/slide\d+\.xml", n)]
                )
                texts: list[str] = []
                for slide_path in slides:
                    xml = z.read(slide_path).decode("utf-8", errors="replace")
                    runs = re.findall(r"<a:t>([^<]*)</a:t>", xml)
                    texts.append(" ".join(runs))
            full = "\n\n".join(texts)
            return full, [{"heading": "", "text": full, "page": None}]
        except Exception:
            return "", []
